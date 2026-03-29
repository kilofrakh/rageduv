# ingest.py
"""
Stage 1 — Document Ingestion Pipeline
--------------------------------------
Handles: PDF, DOCX, PPTX, TXT
Steps:
  1. Parse raw file → plain text + metadata
  2. Split into overlapping chunks
  3. Embed each chunk
  4. Upsert into ChromaDB (deduplicates by file hash)
"""

import hashlib
import os
from pathlib import Path
from typing import List

from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings

import config


# ── Embeddings (loaded once, reused everywhere) ───────────────────────────────
_embeddings = None

def get_embeddings() -> HuggingFaceEmbeddings:
    global _embeddings
    if _embeddings is None:
        _embeddings = HuggingFaceEmbeddings(model_name=config.EMBEDDING_MODEL)
    return _embeddings


# ── Vector store ──────────────────────────────────────────────────────────────
def get_vectorstore() -> Chroma:
    return Chroma(
        collection_name=config.COLLECTION_NAME,
        embedding_function=get_embeddings(),
        persist_directory=config.CHROMA_DIR,
    )


# ── Parsers ───────────────────────────────────────────────────────────────────
def _parse_pdf(path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _parse_docx(path: str) -> str:
    from docx import Document as DocxDocument
    doc = DocxDocument(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides_text)


def _parse_txt(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")


PARSERS = {
    ".pdf":  _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".txt":  _parse_txt,
}


# ── Core ingestion ─────────────────────────────────────────────────────────────
def ingest_file(
    file_path: str,
    subject: str = "General",
    professor: str = "Unknown",
) -> dict:
    """
    Ingest a single file into the vector store.
    Returns a summary dict with chunk count and file metadata.
    """
    ext = Path(file_path).suffix.lower()
    if ext not in PARSERS:
        raise ValueError(f"Unsupported file type: {ext}. Use {list(PARSERS.keys())}")

    # 1. Parse
    raw_text = PARSERS[ext](file_path)
    if not raw_text.strip():
        raise ValueError("File appears to be empty or unreadable.")

    # 2. Deduplication via file hash
    file_hash = hashlib.md5(Path(file_path).read_bytes()).hexdigest()
    filename  = Path(file_path).name

    # 3. Split into chunks
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=config.CHUNK_SIZE,
        chunk_overlap=config.CHUNK_OVERLAP,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    chunks = splitter.split_text(raw_text)

    # 4. Wrap as LangChain Documents with rich metadata
    documents = [
        Document(
            page_content=chunk,
            metadata={
                "source":    filename,
                "file_hash": file_hash,
                "subject":   subject,
                "professor": professor,
                "chunk_idx": i,
            },
        )
        for i, chunk in enumerate(chunks)
    ]

    # 5. Upsert into ChromaDB
    #    Use "<hash>-<chunk_idx>" as stable IDs to avoid duplicates on re-upload
    ids = [f"{file_hash}-{i}" for i in range(len(documents))]
    vs  = get_vectorstore()
    vs.add_documents(documents, ids=ids)

    return {
        "filename":   filename,
        "subject":    subject,
        "chunks":     len(chunks),
        "file_hash":  file_hash,
    }


def ingest_directory(directory: str, subject: str = "General", professor: str = "Unknown") -> List[dict]:
    """Batch ingest every supported file in a folder."""
    results = []
    for f in Path(directory).iterdir():
        if f.suffix.lower() in PARSERS:
            try:
                result = ingest_file(str(f), subject=subject, professor=professor)
                results.append(result)
                print(f"  ✓ {f.name}  ({result['chunks']} chunks)")
            except Exception as e:
                print(f"  ✗ {f.name}: {e}")
    return results


# ── CLI helper ────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python ingest.py <file_or_folder> [subject] [professor]")
        sys.exit(1)
    path    = sys.argv[1]
    subject = sys.argv[2] if len(sys.argv) > 2 else "General"
    prof    = sys.argv[3] if len(sys.argv) > 3 else "Unknown"

    if os.path.isdir(path):
        results = ingest_directory(path, subject, prof)
        print(f"\nIngested {len(results)} files.")
    else:
        result = ingest_file(path, subject, prof)
        print(f"Ingested: {result}")
